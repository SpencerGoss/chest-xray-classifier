"""Builds Chest_Xray_Project.pptx, a tight, ~9-minute deck (12 slides) with a featured
live demo, an EDA slide (condition co-occurrence + demographics), a deployment/monitoring
slide, and speaker notes with timing. Plain language, professional theme, no overlaps.
Reuses figures in slides_assets/ (extracted from the notebook) + one custom AUC chart."""
import os
from PIL import Image
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

A = "slides_assets"

# ---- Theme -----------------------------------------------------------------
INK     = RGBColor(0x22, 0x2A, 0x30)
PRIMARY = RGBColor(0x10, 0x49, 0x5C)
PRIMARY2= RGBColor(0x1C, 0x6E, 0x7A)
ACCENT  = RGBColor(0xE0, 0x7A, 0x33)
GREEN   = RGBColor(0x2E, 0x7D, 0x4F)
RED     = RGBColor(0xB0, 0x3A, 0x2E)
LIGHT   = RGBColor(0xF6, 0xF8, 0xF9)
CARD    = RGBColor(0xFF, 0xFF, 0xFF)
MUTE    = RGBColor(0x6E, 0x78, 0x82)
LINE    = RGBColor(0xDD, 0xE3, 0xE6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PALE    = RGBColor(0xCF, 0xE0, 0xE4)
FONT    = "Calibri"
H_TEAL, H_AMBER, H_GRAY = "#10495C", "#E07A33", "#9AA4AC"

MODELS = [  # name, test_auc (from the shipped Colab run)
    ("Neural network", 0.729), ("Random forest", 0.703), ("Naive Bayes", 0.686),
    ("Logistic regression", 0.677), ("Gradient boosting", 0.672),
    ("k-Nearest Neighbors", 0.610), ("Baseline (says nothing)", 0.500),
]

# Deep-learning progression (test AUC) from the GPU run: from scratch -> transfer -> fine-tuned.
PROGRESSION = [("CNN\nfrom scratch", 0.644), ("Transfer learning\n(frozen)", 0.729),
               ("Fine-tuned\n(unfrozen)", 0.738)]

def chart_progression(path):
    names = [n for n, _ in PROGRESSION]; vals = [v for _, v in PROGRESSION]
    colors = [H_TEAL, H_TEAL, H_AMBER]
    fig, ax = plt.subplots(figsize=(7.6, 4.6))
    bars = ax.bar(names, vals, color=colors, edgecolor="white", width=0.6)
    ax.axhline(0.5, color="#B6BEC4", ls="--", lw=1.2)
    ax.text(2.45, 0.505, "chance", color="#8A949C", fontsize=9, ha="right")
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width()/2, v + 0.006, f"{v:.2f}", ha="center", fontsize=13,
                fontweight="bold", color="#3A444C")
    ax.set_ylim(0.45, 0.80)
    ax.set_ylabel("Test AUC", fontsize=12, color="#3A444C")
    ax.set_title("Deep learning: each step adds accuracy", fontsize=15, fontweight="bold",
                 color="#222A30", pad=12)
    for s in ["top", "right"]: ax.spines[s].set_visible(False)
    ax.tick_params(labelsize=11.5, colors="#3A444C")
    ax.set_axisbelow(True); ax.yaxis.grid(True, color="#EEF1F3")
    fig.tight_layout(); fig.savefig(path, dpi=150, facecolor="white"); plt.close(fig)

def chart_auc(path):
    data = sorted(MODELS, key=lambda x: x[1])
    names = [m for m, _ in data]; vals = [v for _, v in data]
    # Highlight the neural network (the model we keep), not the technical max, since the
    # top few are tied within noise.
    colors = [H_GRAY if "Baseline" in m else (H_AMBER if m == "Neural network" else H_TEAL)
              for m, v in data]
    fig, ax = plt.subplots(figsize=(8.4, 4.5))
    bars = ax.barh(names, vals, color=colors, edgecolor="white", height=0.66)
    ax.axvline(0.5, color="#B6BEC4", ls="--", lw=1.2)
    for b, v in zip(bars, vals):
        ax.text(v + 0.004, b.get_y() + b.get_height()/2, f"{v:.2f}",
                va="center", ha="left", fontsize=11, color="#3A444C")
    ax.set_xlim(0.45, 0.78)
    ax.set_xlabel("AUC, higher is better (0.5 = guessing)", fontsize=11.5, color="#3A444C")
    ax.set_title("How the models compared", fontsize=15, fontweight="bold", color="#222A30", pad=12)
    for s in ["top", "right"]: ax.spines[s].set_visible(False)
    ax.tick_params(labelsize=11, colors="#3A444C")
    ax.set_axisbelow(True); ax.xaxis.grid(True, color="#EEF1F3")
    fig.tight_layout(); fig.savefig(path, dpi=150, facecolor="white"); plt.close(fig)

chart_auc(os.path.join(A, "auc_simple.png"))
chart_progression(os.path.join(A, "progression.png"))

# ---- PPTX helpers ----------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH, MARGIN = 13.333, 7.5, 0.6
BLANK = prs.slide_layouts[6]

def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = bg
    return s

def rect(s, l, t, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space; p.space_after = Pt(4)
        if isinstance(para, tuple): para = [para]
        for (txt, size, color, bold, *rest) in para:
            r = p.add_run(); r.text = txt
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold
            if rest and rest[0]: r.font.italic = True
    return tb

def bullets(s, l, t, w, h, items, size=15, color=INK, gap=8, mark=ACCENT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08
        r1 = p.add_run(); r1.text = "•  "
        r1.font.name = FONT; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = mark
        r2 = p.add_run(); r2.text = it
        r2.font.name = FONT; r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb

def img_fit(s, path, l, t, w, h, frame=True):
    W, H = Image.open(path).size
    ar = W / H; bar = w / h
    if ar > bar: nw, nh = w, w / ar
    else: nh, nw = h, h * ar
    nl, nt = l + (w - nw) / 2, t + (h - nh) / 2
    if frame:
        rect(s, nl - 0.06, nt - 0.06, nw + 0.12, nh + 0.12, CARD, line=LINE)
    s.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))

def header(s, title, num):
    rect(s, 0, 0, 0.22, SH, PRIMARY)
    text(s, MARGIN, 0.42, SW - 2*MARGIN, 0.7, [[(title, 26, PRIMARY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, MARGIN, 1.18, 1.15, 0.07, ACCENT)
    text(s, SW - 2.4, 7.06, 1.9, 0.3, [[(f"{num}", 11, MUTE, False)]], align=PP_ALIGN.RIGHT)
    text(s, MARGIN, 7.06, 8.0, 0.3, [[("NIH ChestX-ray14  ·  MSBA 645 Final Project", 10, MUTE, False)]])

def stat(s, l, t, w, h, big, lab, color=PRIMARY):
    rect(s, l, t, w, h, CARD, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, l, t, w, 0.12, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, l, t + 0.18, w, h - 0.26, [[(big, 28, color, True)], [(lab, 12.5, MUTE, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

# ===================================================================== SLIDES
# 1  Title
s = slide(PRIMARY)
rect(s, 0, 0, 0.35, SH, ACCENT)
text(s, 1.0, 2.1, 11.3, 2.0,
     [[("Detecting Chest Conditions", 46, WHITE, True)], [("from X-rays", 46, WHITE, True)]])
rect(s, 1.05, 4.05, 2.2, 0.08, ACCENT)
text(s, 1.0, 4.3, 11.3, 1.0,
     [[("Predicting 14 conditions from a chest X-ray with deep learning", 19, PALE, False)]])
text(s, 1.0, 6.3, 11.3, 0.9,
     [[("MSBA 645 Final Project", 14, WHITE, True)], [("Team: [add names here]", 13, PALE, False)]])
notes(s, "(~15 seconds. Say this:)  Good morning. For our project, we built a tool that looks at a "
         "chest X-ray and predicts which of 14 different chest conditions might be present, things like "
         "pneumonia, fluid around the lungs, or an enlarged heart. We used a large public collection of "
         "real chest X-rays called NIH ChestX-ray14. Over the next few minutes we'll show you the data, "
         "how the model works, how well it does, and then we'll run it live. I'm [name] and this is "
         "[name].")

# 2  Problem + data
s = slide(); header(s, "The problem, and the data", 2)
text(s, MARGIN, 1.6, SW - 2*MARGIN, 0.9, [
    [("Given one chest X-ray, predict which of 14 conditions are present. An X-ray can have several at "
      "once, so the model gives a yes/no for each one.", 16.5, INK, False)]])
cw = (SW - 2*MARGIN - 2*0.3) / 3
for i, (b, l) in enumerate([("112,120", "chest X-rays"), ("30,805", "patients"), ("14", "conditions")]):
    stat(s, MARGIN + i*(cw+0.3), 2.55, cw, 1.15, b, l, color=PRIMARY if i != 2 else ACCENT)
img_fit(s, os.path.join(A, "samples.png"), MARGIN, 3.9, SW - 2*MARGIN, 2.45)
text(s, MARGIN, 6.42, SW - 2*MARGIN, 0.4,
     [[("Real examples with their labels. For this project we used about 25,000 of the 112,120 images "
        "(a laptop-friendly slice).", 12, MUTE, False, True)]],
     align=PP_ALIGN.CENTER)
notes(s, "(~1 minute. Say this:)  Here's the problem we set out to solve. You give the model one "
         "chest X-ray, and it gives back a yes or no for each of 14 conditions. The important part is "
         "that a single X-ray can have more than one thing wrong with it at the same time, so the model "
         "isn't picking just one answer, it's judging all 14 at once. The full dataset is huge: about "
         "112,000 X-rays from around 30,000 patients, each labeled for these 14 conditions. We didn't "
         "use all of it. We used about 25,000 images, split into a training set, a validation set, and "
         "a test set, so the whole project runs on a normal laptop. On the screen are a few real "
         "examples with their actual labels, so you can see what the model is working with.")

# 3  Imbalance / why accuracy lies
s = slide(); header(s, "The catch: the data is lopsided", 3)
img_fit(s, os.path.join(A, "freq.png"), MARGIN, 1.65, 7.0, 3.9)
rect(s, 7.8, 1.9, SW - MARGIN - 7.8, 3.5, RGBColor(0xFB, 0xEE, 0xE2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 8.05, 2.15, SW - MARGIN - 8.3, 3.0, [
    [("Why this matters", 15, ACCENT, True)],
    [("About half of all X-rays are normal, and some conditions are rare.", 14.5, INK, False)],
    [(" ", 7, INK, False)],
    [("So a lazy model that always says “nothing’s wrong” is right ~95% of the time, but it’s useless.",
      14.5, INK, False)],
    [(" ", 7, INK, False)],
    [("That’s why we don’t trust accuracy alone. We judge by AUC (how well it tells sick from healthy) "
      "and recall (how much real disease it catches).", 14.5, INK, False)],
], space=1.05)
notes(s, "(~1 minute. Say this:)  Now the catch, and this is the most important idea in our whole "
         "project. The data is very lopsided. About half of all the X-rays are completely normal, and "
         "several of the conditions are quite rare. So imagine a lazy model that just says 'nothing is "
         "wrong' on every single image. Because most images really are normal, that model would be "
         "right about 95 percent of the time. Ninety-five percent sounds great, but it caught zero "
         "actual disease, so it's useless. That's why, for the rest of this talk, we don't judge our "
         "model on accuracy. We use two better measures instead: AUC, which is how well the model "
         "separates sick patients from healthy ones, and recall, which is how much of the real disease "
         "it actually catches.")

# 4  A closer look at the data (EDA: co-occurrence + demographics)
s = slide(); header(s, "A closer look at the data", 4)
img_fit(s, os.path.join(A, "cooccurrence.png"), MARGIN, 1.55, 6.0, 4.1)
text(s, MARGIN, 5.72, 6.0, 0.4,
     [[("How often two conditions appear on the same X-ray (brighter = more often together).",
        11, MUTE, False, True)]], align=PP_ALIGN.CENTER)
text(s, 7.0, 1.65, SW - MARGIN - 7.0, 0.5, [[("Conditions overlap", 15, PRIMARY, True)]])
bullets(s, 7.0, 2.2, SW - MARGIN - 7.0, 2.0, [
    "Findings travel together: Effusion + Infiltration, Atelectasis + Effusion, Pneumothorax + Emphysema.",
    "So the 14 labels aren’t independent, one image can carry several at once.",
    "That overlap is why we score each condition on its own with AUC and recall, not one accuracy number.",
], size=13.5)
text(s, 7.0, 4.35, SW - MARGIN - 7.0, 0.5, [[("Who’s in the data", 15, PRIMARY, True)]])
img_fit(s, os.path.join(A, "demographics.png"), 7.0, 4.8, SW - MARGIN - 7.0, 1.9, frame=False)
text(s, 7.0, 6.55, SW - MARGIN - 7.0, 0.4,
     [[("Mostly middle-aged adults (median 49), 56% male / 44% female, a generalisation flag.",
        11, MUTE, False, True)]], align=PP_ALIGN.CENTER)
notes(s, "(~45 seconds. Say this:)  Before building anything, we looked closely at the data and "
         "found two things worth sharing. First, this heat map shows that the conditions often appear "
         "together on the same X-ray. The brighter squares are the pairs that show up together a lot, "
         "for example fluid around the lung often comes with infiltration. This matters because it "
         "means the 14 conditions aren't separate, independent things, which is another reason a single "
         "accuracy score doesn't fit here. Second, we looked at who's actually in the data. Most of the "
         "patients are middle-aged adults, the median age is about 49, and there are a few more men "
         "than women. We point this out because a model only really learns the kind of patients it was "
         "trained on, so that's something to keep in mind before trusting it on a different group.")

# 5  Approach
s = slide(); header(s, "Our approach: borrow a trained network", 5)
text(s, MARGIN, 1.6, SW - 2*MARGIN, 0.9, [
    [("Training an image model from scratch needs a GPU and tons of data. Instead we reuse a network "
      "already trained on millions of images, and train small models on top of it.", 16, INK, False)]])
steps = ["Chest\nX-ray", "Pretrained network\n(frozen)", "A list of\nnumbers", "Train simple\nmodels",
         "14 condition\nprobabilities"]
bw, bh, bt = 2.05, 1.5, 3.7
gap = (SW - 2*MARGIN - 5*bw) / 4
cols = [PRIMARY, PRIMARY2, PRIMARY2, PRIMARY2, ACCENT]
for i, (st, col) in enumerate(zip(steps, cols)):
    l = MARGIN + i*(bw+gap)
    rect(s, l, bt, bw, bh, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, l, bt, bw, bh, [[(st, 14.5, WHITE, True)]], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    if i < 4:
        rect(s, l + bw + gap/2 - 0.18, bt + bh/2 - 0.18, 0.36, 0.36, ACCENT, shape=MSO_SHAPE.CHEVRON)
text(s, MARGIN, 5.55, SW - 2*MARGIN, 0.9,
     [[("This is called transfer learning. It’s one of several deep-learning techniques we use:", 14, MUTE, False, True)],
      [("transfer learning, a CNN built from scratch, fine-tuning, dropout, and Grad-CAM.", 14, MUTE, False, True)]])
notes(s, "(~45 seconds. Say this:)  Here's how the model actually works. Training an image model "
         "completely from scratch needs a massive amount of data and a powerful graphics card, which we "
         "didn't have. So instead, we borrowed a network called EfficientNet that was already trained "
         "on millions of everyday photos. It already knows how to pick out edges, shapes, and textures. "
         "We freeze it, meaning we don't change it at all, and just use it to turn each X-ray into a "
         "list of 1,280 numbers. Then we train small, simple models on those numbers to predict the 14 "
         "conditions. Following the arrows: the X-ray goes in, the pretrained network turns it into "
         "numbers, and our model turns those numbers into a probability for each condition. This is "
         "called transfer learning, and it's what let us do this without a fancy GPU.")

# 6  Model comparison
s = slide(); header(s, "We tried several models", 6)
img_fit(s, os.path.join(A, "auc_simple.png"), MARGIN, 1.65, 7.7, 4.9)
text(s, 8.6, 2.0, SW - MARGIN - 8.6, 0.5, [[("What we found", 15, PRIMARY, True)]])
bullets(s, 8.6, 2.5, SW - MARGIN - 8.6, 2.3, [
    "Our deep-learning model (a neural network) came out best, around 0.73.",
    "The classical models (logistic regression, random forest, etc.) are simpler baselines for comparison.",
    "It's the one we keep, tune, and build on in the next slides.",
], size=14)
rect(s, 8.6, 4.75, SW - MARGIN - 8.6, 1.95, RGBColor(0xFB, 0xEE, 0xE2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 8.85, 4.95, SW - MARGIN - 9.1, 1.6, [
    [("The catch: ", 13, ACCENT, True),
     ("random forest came close on AUC, but it memorized the training data (perfect on train, 0.70 on "
      "new data) and rarely says “yes,” so it misses real cases. A high AUC isn’t enough, so we also "
      "check recall.", 13, INK, False)],
], space=1.04)
notes(s, "(~1 minute. Say this:)  We didn't just build one model and hope for the best. We tried "
         "several different models on exactly the same data and compared them fairly. The bar at the "
         "very bottom is a baseline that's basically guessing, it sits at 0.50. Our neural network came "
         "out on top, at about 0.73, and that's the one we kept. One thing worth pointing out: the "
         "random forest looked almost as good on this chart, but it had basically memorized the "
         "training data, it scored perfectly on data it had already seen and much worse on new data, "
         "and it almost never actually flags a condition, so it would miss real cases. That's a great "
         "example of why a high score on this chart isn't the whole story, and why we also care about "
         "recall, which is catching real disease.")

# 7  Deep learning progression
s = slide(); header(s, "Deep learning: from scratch to fine-tuned", 7)
img_fit(s, os.path.join(A, "progression.png"), MARGIN, 1.6, 7.4, 5.0)
text(s, 8.3, 2.0, SW - MARGIN - 8.3, 0.5, [[("The deep-learning story", 15, PRIMARY, True)]])
bullets(s, 8.3, 2.55, SW - MARGIN - 8.3, 3.6, [
    "A CNN built from scratch reaches only 0.64, with limited data it can't learn good features on its own.",
    "Transfer learning (reusing a pretrained network) jumps to 0.73, the big win.",
    "Fine-tuning the whole network nudges it to 0.74.",
], size=14)
rect(s, 8.3, 5.2, SW - MARGIN - 8.3, 1.35, RGBColor(0xEC, 0xF3, 0xF2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 8.55, 5.35, SW - MARGIN - 8.8, 1.1,
     [[("This is the core deep-learning lesson: with limited data, standing on a pretrained network "
        "beats building one from scratch.", 13, INK, False)]], space=1.05)
notes(s, "(~45 seconds. Say this:)  This slide is really the heart of the deep-learning side. We "
         "tried three different deep-learning approaches, and the numbers tell a clean story. A network "
         "built completely from scratch only reached 0.64, because with our limited data it just can't "
         "learn good features on its own. When we reused the pretrained network instead, that's "
         "transfer learning, the score jumped to 0.73, and that's the big win. And when we fine-tuned, "
         "letting the network adjust just a little to chest X-rays specifically, it went up a bit more "
         "to 0.74. So the takeaway is simple: when you don't have a ton of data, standing on a network "
         "that's already been trained beats building one from scratch.")

# 8  Best model results
s = slide(); header(s, "Our best model: how good is it?", 8)
text(s, MARGIN, 1.6, 5.2, 0.5, [[("The neural network", 16, PRIMARY, True)]])
stat(s, MARGIN, 2.15, 2.4, 1.2, "0.74", "Accuracy", color=PRIMARY)
stat(s, MARGIN + 2.7, 2.15, 2.4, 1.2, "0.73", "AUC", color=ACCENT)
text(s, MARGIN, 3.65, 5.2, 2.6, [
    [("Train vs. test", 14, PRIMARY, True)],
    [("Accuracy:  0.75 on training,  0.74 on new data", 14.5, INK, False)],
    [("AUC:  0.84 on training,  0.73 on new data", 14.5, INK, False)],
    [(" ", 7, INK, False)],
    [("The test numbers are close to training, so it’s genuinely learning, not just memorizing.",
      13.5, MUTE, False, True)],
], space=1.1)
img_fit(s, os.path.join(A, "confusion.png"), 6.3, 1.55, SW - MARGIN - 6.3, 5.2)
notes(s, "(~1 minute. Say this:)  So how good is our chosen model? On X-rays it had never seen "
         "before, it gets about 74 percent accuracy and an AUC of 0.73. Just as important, the scores "
         "on the training data and on new data are close together, 0.84 versus 0.73 on AUC, which tells "
         "us the model is genuinely learning the patterns rather than just memorizing examples. On the "
         "right is what's called a confusion matrix. I won't go through every box, but the strong line "
         "down the diagonal is where the model got it right, and the cells off that line are the "
         "mix-ups, usually between conditions that look alike on an X-ray. With that, let's actually see "
         "it in action.")

# 7  DEMO
s = slide(PRIMARY)
rect(s, 0, 0, 0.35, SH, ACCENT)
text(s, MARGIN, 0.55, SW - 2*MARGIN, 0.9, [[("Live demo", 32, WHITE, True)]])
rect(s, MARGIN, 1.5, 1.3, 0.08, ACCENT)
text(s, MARGIN, 1.75, 5.7, 4.6, [
    [("We’ll open the app and try it live:", 17, WHITE, True)],
    [(" ", 8, WHITE, False)],
])
bullets(s, MARGIN, 2.5, 5.6, 3.6, [
    "Upload a chest X-ray.",
    "See the probability for each of the 14 conditions.",
    "See which conditions it flags.",
    "Look at the Grad-CAM heat map, where the model is looking.",
    "Show one it gets right, and one it gets wrong.",
    "Then switch to the Monitoring tab (next slide).",
], size=15.5, color=WHITE, mark=ACCENT)
img_fit(s, os.path.join(A, "gradcam_correct.png"), 6.7, 1.75, SW - MARGIN - 6.7, 2.25, frame=False)
img_fit(s, os.path.join(A, "gradcam_wrong.png"), 6.7, 4.15, SW - MARGIN - 6.7, 2.25, frame=False)
text(s, 6.7, 6.45, SW - MARGIN - 6.7, 0.4,
     [[("Heat maps from the app (backup if the live demo can’t run).", 11.5, PALE, False, True)]],
     align=PP_ALIGN.CENTER)
notes(s, "(~2 to 3 minutes. Have the app already open before this slide. Say this as you click "
         "through:)  This is the live demo. Let me pick a chest X-ray. Here you can see the model's "
         "probability for each of the 14 conditions, and over here it flags the ones it thinks are "
         "present. You might notice the flagged condition isn't always the one with the highest number, "
         "and that's on purpose: instead of one fixed cutoff, each condition has its own threshold that "
         "we tuned to favor catching disease. Now here's the interesting part. This heat map, called "
         "Grad-CAM, shows where the model was actually looking when it made its decision. You can see "
         "it's focused on the lungs, not the background, which gives us some confidence it's reasoning "
         "sensibly. Let me show you one the model gets right, and then one it gets wrong, so you can see "
         "both.  (If the live app won't run, the two heat maps on this slide are your backup.)")

# 9  From demo to product: monitoring
s = slide(); header(s, "From demo to product: monitoring", 9)
text(s, MARGIN, 1.55, SW - 2*MARGIN, 0.85, [
    [("The app doesn’t just predict, it keeps a record. Every classification is logged, so you can watch "
      "the model across many cases, not just one image.", 16, INK, False)]])
steps = ["Classify\nan X-ray", "Log every\nprediction (CSV)", "Monitoring\ndashboard",
         "or Power BI /\nSheets"]
bw, bh, bt = 2.55, 1.35, 2.65
gap = (SW - 2*MARGIN - 4*bw) / 3
cols = [PRIMARY, PRIMARY2, PRIMARY2, ACCENT]
for i, (stp, col) in enumerate(zip(steps, cols)):
    l = MARGIN + i*(bw+gap)
    rect(s, l, bt, bw, bh, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, l, bt, bw, bh, [[(stp, 14.5, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        rect(s, l + bw + gap/2 - 0.18, bt + bh/2 - 0.18, 0.36, 0.36, ACCENT, shape=MSO_SHAPE.CHEVRON)
text(s, MARGIN, 4.4, SW - 2*MARGIN, 0.5, [[("What the dashboard tracks", 15, PRIMARY, True)]])
colw = (SW - 2*MARGIN - 0.4) / 2
bullets(s, MARGIN, 4.95, colw - 0.2, 1.6, [
    "How often each of the 14 conditions gets flagged.",
    "The spread of the model’s confidence.",
], size=13.5)
bullets(s, MARGIN + colw + 0.4, 4.95, colw - 0.2, 1.6, [
    "Most recent cases, at a glance.",
    "User feedback: a “looks right / looks off” tap on each prediction.",
], size=13.5)
rect(s, MARGIN, 6.35, SW - 2*MARGIN, 0.7, RGBColor(0xEC, 0xF3, 0xF2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, MARGIN + 0.3, 6.35, SW - 2*MARGIN - 0.6, 0.7, [
    [("Why it matters:  ", 13, PRIMARY, True),
     ("one prediction is a demo; logging turns it into something you can actually monitor in use. The "
      "same CSV is the single source any BI tool can read.", 13, INK, False)]],
    anchor=MSO_ANCHOR.MIDDLE, space=1.03)
notes(s, "(~45 seconds. Say this:)  The app doesn't just make a prediction and forget it. Every "
         "single time it classifies an X-ray, it saves a record: the image, the top prediction, all 14 "
         "probabilities, and how many conditions were flagged. This Monitoring tab reads all those "
         "records and gives a big-picture view: how often each condition gets flagged, how confident "
         "the model tends to be, the most recent cases, and any thumbs-up or thumbs-down feedback that "
         "users gave. The reason this matters is that a single prediction is just a demo, but keeping a "
         "log is what lets you actually watch how the model behaves over time in real use. And because "
         "it's just a spreadsheet file, that same data could feed a business tool like Power BI.")

# 10  Lessons + state of the art
s = slide(); header(s, "What we learned", 10)
colw = (SW - 2*MARGIN - 0.4) / 2
for i, (htext, hcol, items) in enumerate([
    ("What worked", GREEN, [
        "Telling the model to pay attention to rare cases (instead of ignoring them).",
        "Preparing the images the way the network expects.",
        "Checking AUC and recall, not just accuracy.",
        "Grad-CAM showed it looks at the lungs, not the background."]),
    ("What we tested / next", PRIMARY, [
        "Labels are only about 90% accurate (auto-pulled from reports), so some are simply wrong.",
        "We tried a patient-level split (no patient on both sides): AUC 0.67 vs 0.73, so random splitting was a bit optimistic.",
        "We checked if more data helps: still rising slowly, so a little, but the bigger win is fine-tuning.",
        "The deep-learning progression (CNN 0.64 -> transfer 0.73 -> fine-tuned 0.74) is the main story."]),
]):
    l = MARGIN + i*(colw+0.4)
    rect(s, l, 1.6, colw, 0.55, hcol, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, l, 1.6, colw, 0.55, [[(htext, 15, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, l + 0.15, 2.4, colw - 0.3, 3.0, items, size=13.5, gap=7, mark=hcol)
rect(s, MARGIN, 5.75, SW - 2*MARGIN, 0.95, RGBColor(0xEC, 0xF3, 0xF2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, MARGIN + 0.3, 5.75, SW - 2*MARGIN - 0.6, 0.95, [
    [("How this compares:", 14, PRIMARY, True),
     ("  the best published model (CheXNet) hits 0.7 to 0.85 AUC, but it’s fully trained on all 112,000",
      14, INK, False)],
    [("images with a GPU. Ours is a lighter, laptop-friendly version of the same idea, 0.73 frozen and 0.74 fine-tuned.",
      14, INK, False)],
], anchor=MSO_ANCHOR.MIDDLE, space=1.05)
notes(s, "(~1 minute. Say this:)  A few things we learned. On what worked: telling the model to pay "
         "extra attention to the rare conditions instead of ignoring them, preparing the images exactly "
         "the way the network expects, judging the model on AUC and recall instead of accuracy, and "
         "using the heat maps to confirm it looks at the lungs and not the background. We also want to "
         "be honest about the limits. The labels in this dataset were pulled automatically from "
         "doctors' reports and are only about 90 percent accurate, so some of them are simply wrong, "
         "which caps how good any model can look. We also tried a stricter test where no patient "
         "appears in both training and testing, and the score dropped from 0.73 to 0.67, so our main "
         "number is a little optimistic. For comparison, the best published model in this area scores "
         "between 0.7 and 0.85, but it uses all 112,000 images and a powerful GPU. Ours is a much "
         "lighter, laptop-friendly version, and it lands in the same ballpark.")

# 11  Thank you
s = slide(PRIMARY)
rect(s, 0, 0, 0.35, SH, ACCENT)
text(s, 1.0, 2.5, 11.3, 1.4, [[("Thank you", 48, WHITE, True)]])
rect(s, 1.05, 4.0, 2.2, 0.08, ACCENT)
text(s, 1.0, 4.3, 11.3, 1.0, [[("Questions?", 22, PALE, False)]])
text(s, 1.0, 6.4, 11.3, 0.6,
     [[("Chest X-ray multi-label classifier  ·  NIH ChestX-ray14  ·  MSBA 645", 13, PALE, False)]])
notes(s, "(~15 seconds. Say this:)  So to sum up: we built a tool that predicts 14 chest conditions "
         "from an X-ray using transfer learning, we measured it honestly with AUC and recall instead of "
         "accuracy, and we wrapped it in a working app that explains its reasoning and tracks its own "
         "usage. Thank you for listening, and we're happy to take any questions.")

out = "Chest_Xray_Project.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides.")
